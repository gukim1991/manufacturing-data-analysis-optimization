from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0xE8, 0x7A, 0x2B)
GRAY_BG = RGBColor(0xF4, 0xF6, 0xF8)
DARK = RGBColor(0x22, 0x2B, 0x3A)
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x1F, 0x77, 0xB4)


def add_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullet_box(slide, left, top, width, height, items, font_size=14, color=DARK, font_name='맑은 고딕'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tb


def add_header_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(11), Inches(0.55),
                 title, font_size=24, bold=True, color=LIGHT)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    if subtitle:
        add_text_box(slide, Inches(0.4), Inches(1.0), Inches(12), Inches(0.4),
                     subtitle, font_size=13, color=DARK)


def add_card(slide, left, top, width, height, title, body_lines, title_color=NAVY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = NAVY
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                 title, font_size=14, bold=True, color=title_color)
    add_bullet_box(slide, left + Inches(0.15), top + Inches(0.55),
                   width - Inches(0.3), height - Inches(0.6),
                   body_lines, font_size=11, color=DARK)


# ========== Slide 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.08))
deco.fill.solid(); deco.fill.fore_color.rgb = ACCENT; deco.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.0),
             "MIMII 데이터셋 기반 산업 밸브 음향 이상탐지",
             font_size=40, bold=True, color=LIGHT)
add_text_box(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.6),
             "CNN 오토인코더 모델의 구조와 취약점 분석",
             font_size=22, color=ACCENT)
add_text_box(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.5),
             "제조 데이터 분석과 최적화  |  7주차 결과보고서",
             font_size=16, color=LIGHT)
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             "충북대학교 대학원  |  2026.04.28",
             font_size=12, color=LIGHT)

# ========== Slide 2: Agenda ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "목  차  (Agenda)")
agenda = [
    "1.  프로젝트 개요 및 목적",
    "2.  파이프라인 구조 (Step 1 ~ Step 4)",
    "3.  CNN 오토인코더 아키텍처",
    "4.  CNN 오토인코더 모델의 취약점 (10가지)",
    "5.  개선 방향 및 결론",
]
add_bullet_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5),
               agenda, font_size=22, color=NAVY)

# ========== Slide 3: Project Overview ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "1. 프로젝트 개요",
                "MIMII 데이터셋(0dB, valve id_02) 정상 작동음 학습 → 복원오차 기반 비지도 이상탐지")

add_card(slide, Inches(0.4), Inches(1.7), Inches(4.1), Inches(2.5),
         "데이터셋",
         ["MIMII (Malfunctioning Industrial Machine)",
          "장비 종류: Valve (밸브)",
          "기기 ID: id_02",
          "환경 노이즈: 0 dB SNR",
          "샘플링: 16,000 Hz, 10초 단위"])

add_card(slide, Inches(4.6), Inches(1.7), Inches(4.1), Inches(2.5),
         "접근 방법",
         ["정상음만 학습 (Unsupervised)",
          "Mel-Spectrogram → 2D 이미지 변환",
          "CNN AutoEncoder로 복원",
          "복원오차(MSE) 임계값 비교",
          "초과 시 ABNORMAL로 판정"])

add_card(slide, Inches(8.8), Inches(1.7), Inches(4.1), Inches(2.5),
         "기대 효과",
         ["비정상 데이터 부족 환경 대응",
          "결함 라벨링 비용 절감",
          "실시간 모니터링 가능",
          "다양한 결함 유형 일반 탐지",
          "예지보전(PdM) 기반 마련"])

add_text_box(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.4),
             "■ 핵심 가설", font_size=16, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.4), Inches(4.9), Inches(12.5), Inches(2.2),
             "정상 데이터로만 학습된 오토인코더는 정상 패턴을 잘 복원하지만, "
             "학습되지 않은 비정상 패턴을 입력하면 복원에 실패하여 MSE가 커진다. "
             "이 차이를 이용해 라벨 없이도 결함을 탐지할 수 있다.",
             font_size=14, color=DARK)

# ========== Slide 4: Pipeline ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "2. 파이프라인 구조",
                "Step 1 (EDA) → Step 2 (학습) → Step 3 (평가) → Step 4 (실시간 추론)")

steps = [
    ("Step 1", "EDA", "Waveform / STFT /\nMel-Spectrogram\n시각화·청음", BLUE),
    ("Step 2", "Training", "정상음만으로\nCNN AE 학습\n50 epoch, MSE", NAVY),
    ("Step 3", "Evaluation", "정상/비정상 50개\n복원오차 분포\nROC·F1 임계값", ACCENT),
    ("Step 4", "Inference", "2초 슬라이딩 윈도우\n실시간 모니터링\n임계값 0.0016", RED),
]

x_start = Inches(0.5)
gap = Inches(0.25)
box_w = Inches(2.85)
box_h = Inches(2.3)

for i, (label, name, desc, color) in enumerate(steps):
    left = Inches(0.5 + i * (2.85 + 0.25))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), box_w, box_h)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    add_text_box(slide, left, Inches(2.15), box_w, Inches(0.4),
                 label, font_size=14, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.55), box_w, Inches(0.5),
                 name, font_size=20, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(3.15), box_w, Inches(1.2),
                 desc, font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow_left = left + box_w + Inches(0.0)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(3.0),
                                        Inches(0.25), Inches(0.5))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = NAVY
        arrow.line.fill.background()

add_text_box(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
             "■ 데이터 흐름", font_size=16, bold=True, color=ACCENT)
add_bullet_box(slide, Inches(0.7), Inches(5.1), Inches(12), Inches(2.0),
               ["입력: 16kHz WAV → 128-band Mel-Spectrogram (n_mels=128)",
                "전처리: power_to_db(ref=np.max) → (mel_db + 80) / 80 정규화 → 128프레임으로 자르기/패딩",
                "텐서 형태: (1, 128, 128) — 단일 채널 2D 이미지",
                "출력: Sigmoid로 0~1 범위 복원 이미지 → MSE 오차 계산"],
               font_size=13, color=DARK)

# ========== Slide 5: CNN AutoEncoder Architecture ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "3. CNN 오토인코더 아키텍처",
                "5층 Conv 인코더 → 64차원 잠재공간 → 5층 ConvTranspose 디코더")

# Encoder block
enc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.8),
                              Inches(5.0), Inches(4.8))
enc.fill.solid(); enc.fill.fore_color.rgb = BLUE; enc.line.fill.background()
add_text_box(slide, Inches(0.4), Inches(1.95), Inches(5.0), Inches(0.5),
             "ENCODER", font_size=18, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
enc_layers = [
    "Input  (1, 128, 128)",
    "Conv2d(1→16,  s=2)  + ReLU + BN",
    "Conv2d(16→32, s=2)  + ReLU + BN",
    "Conv2d(32→64, s=2)  + ReLU + BN",
    "Conv2d(64→128, s=2) + ReLU + BN",
    "Conv2d(128→256, s=2) + ReLU",
    "Flatten → Linear(4096 → 64)",
]
add_bullet_box(slide, Inches(0.6), Inches(2.5), Inches(4.8), Inches(4.0),
               enc_layers, font_size=13, color=LIGHT)

# Latent
latent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(3.4),
                                 Inches(2.1), Inches(1.6))
latent.fill.solid(); latent.fill.fore_color.rgb = ACCENT; latent.line.fill.background()
add_text_box(slide, Inches(5.6), Inches(3.55), Inches(2.1), Inches(0.5),
             "Latent", font_size=15, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.6), Inches(3.95), Inches(2.1), Inches(0.5),
             "z ∈ ℝ⁶⁴", font_size=18, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.6), Inches(4.4), Inches(2.1), Inches(0.4),
             "(64-dim)", font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# Decoder
dec = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.8),
                              Inches(5.0), Inches(4.8))
dec.fill.solid(); dec.fill.fore_color.rgb = NAVY; dec.line.fill.background()
add_text_box(slide, Inches(7.9), Inches(1.95), Inches(5.0), Inches(0.5),
             "DECODER", font_size=18, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
dec_layers = [
    "Linear(64 → 4096) + ReLU → reshape",
    "ConvT(256→128, s=2) + ReLU + BN",
    "ConvT(128→64,  s=2) + ReLU + BN",
    "ConvT(64→32,   s=2) + ReLU + BN",
    "ConvT(32→16,   s=2) + ReLU + BN",
    "ConvT(16→1,    s=2) + Sigmoid",
    "Output (1, 128, 128)",
]
add_bullet_box(slide, Inches(8.1), Inches(2.5), Inches(4.8), Inches(4.0),
               dec_layers, font_size=13, color=LIGHT)

# Arrows
a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.4), Inches(4.0),
                              Inches(0.25), Inches(0.4))
a1.fill.solid(); a1.fill.fore_color.rgb = DARK; a1.line.fill.background()
a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.7), Inches(4.0),
                              Inches(0.25), Inches(0.4))
a2.fill.solid(); a2.fill.fore_color.rgb = DARK; a2.line.fill.background()

add_text_box(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5),
             "▶ 손실 함수: MSELoss   |   최적화: Adam(lr=0.001)   |   학습 epoch: 50   |   "
             "출력 정규화: Sigmoid (0~1)",
             font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ========== Slide 6: Vulnerabilities Overview ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "4. CNN 오토인코더 모델의 취약점",
                "본 프로젝트(MIMII valve id_02 / 0dB)에서 식별된 10가지 핵심 약점")

vuln_titles = [
    ("①", "도메인 편향", "id_02 단일 밸브만 학습"),
    ("②", "Generalization Gap", "비정상도 잘 복원해버림"),
    ("③", "시간축 손실", "앞 128프레임만 사용"),
    ("④", "정규화 취약성", "ref=np.max 상대 dB"),
    ("⑤", "단일 임계값", "0.0016 하드코딩"),
    ("⑥", "MSE 손실 한계", "국소·희소 결함 묻힘"),
    ("⑦", "강건성 부재", "노이즈 증강 없음"),
    ("⑧", "BatchNorm 의존", "분포 변화에 민감"),
    ("⑨", "위치 비민감성", "CNN translation invariance"),
    ("⑩", "학습-추론 불일치", "10초 학습 vs 2초 윈도우"),
]

cols = 5
rows = 2
cw = Inches(2.45)
ch = Inches(2.3)
xs = Inches(0.45)
ys = Inches(1.95)
gx = Inches(0.13)
gy = Inches(0.2)

for idx, (num, title, sub) in enumerate(vuln_titles):
    r = idx // cols
    c = idx % cols
    left = Inches(0.45 + c * (2.45 + 0.13))
    top = Inches(1.95 + r * (2.3 + 0.2))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = NAVY; card.line.width = Pt(0.75)
    card.shadow.inherit = False

    num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15),
                                       Inches(0.55), Inches(0.55))
    num_box.fill.solid(); num_box.fill.fore_color.rgb = RED; num_box.line.fill.background()
    add_text_box(slide, left + Inches(0.15), top + Inches(0.18), Inches(0.55), Inches(0.5),
                 num, font_size=18, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.2), cw - Inches(0.9), Inches(0.5),
                 title, font_size=14, bold=True, color=NAVY)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.95), cw - Inches(0.4), ch - Inches(1.1),
                 sub, font_size=12, color=DARK)

# ========== Slide 7: Vulnerabilities Detail 1 (1~5) ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "4-1. 취약점 상세 (1) — 데이터·전처리·임계값",
                "학습 데이터 범위와 전처리 파이프라인에서 발생하는 5가지 약점")

details_1 = [
    ("① 학습 데이터 도메인 편향",
     "step2_train_audio_ae.py:15  →  id_02/normal 한 폴더만 사용. 다른 ID(00, 04, 06)나 SNR(6dB, -6dB)에서는 정상음조차 복원오차가 커져 오탐(False Alarm) 폭증."),
    ("② Generalization Gap (과도한 복원 능력)",
     "5층 ConvTranspose 디코더의 표현력이 강해, 학습하지 않은 비정상 패턴까지 깔끔하게 복원해버림. 정상/비정상 MSE 분포가 겹쳐 임계값 분리가 어려워짐."),
    ("③ 시간축 정보 손실",
     "step2_train_audio_ae.py:50  →  mel_db[:, :128]. 10초 오디오 중 약 4초만 사용. 결함 신호가 후반부에만 나타나면 학습/평가 모두 놓침."),
    ("④ 정규화 방식 취약성",
     "power_to_db(mel, ref=np.max)는 파일별 최댓값 기준 상대 dB. 큰 충격음이 들어오면 max가 올라가 전체 스펙트로그램이 시프트 → 음량 변화에 비강건."),
    ("⑤ 단일 고정 임계값 (Threshold)",
     "step4_inference_audio_ae.py:70  →  THRESHOLD = 0.0016 하드코딩. 환경 노이즈·온도·마이크 위치 변화에 따른 분포 이동(Concept Drift)에 무력."),
]

y = Inches(1.7)
for title, body in details_1:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(0.12), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text_box(slide, Inches(0.65), y, Inches(12.3), Inches(0.4),
                 title, font_size=15, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.65), y + Inches(0.4), Inches(12.3), Inches(0.65),
                 body, font_size=11, color=DARK)
    y += Inches(1.1)

# ========== Slide 8: Vulnerabilities Detail 2 (6~10) ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "4-2. 취약점 상세 (2) — 손실·구조·추론",
                "손실 함수 / 모델 구조 / 학습-추론 정합성에서 발생하는 5가지 약점")

details_2 = [
    ("⑥ MSE 손실의 한계 — 국소 결함 묻힘",
     "MSE는 16,384픽셀(128×128) 평균. 결함이 특정 주파수 대역의 짧은 클릭처럼 국소·희소(sparse)하면 평균에 묻혀 정상 수준 MSE가 나옴. SSIM이나 patch-wise max error가 더 적합."),
    ("⑦ 적대적/노이즈 강건성 부재",
     "0dB 배경 노이즈가 학습에 그대로 포함. 추론 시 다른 종류의 노이즈(공장 교체, 새 컴프레서 가동)가 들어오면 학습 분포 밖이라 정상도 비정상으로 판정. 데이터 증강(SpecAugment) 부재."),
    ("⑧ BatchNorm 의존성",
     "인코더에 BatchNorm2d 다수 포함. eval()에서 running stats 사용하나, 이는 id_02 정상 데이터의 통계라 분포 변화에 민감. LayerNorm/InstanceNorm이 더 안정적."),
    ("⑨ CNN의 위치 비민감성 (Translation Invariance)",
     "Conv는 위치에 둔감해서 결함 발생 위치 정보를 잠재공간 64차원에 잘 못 담음. 시간축이 약간만 어긋난 정상음도 비정상처럼 복원될 수 있음."),
    ("⑩ 학습-추론 분포 불일치",
     "step4 슬라이딩 윈도우는 2초/1초hop이지만, 학습은 10초 전체에서 앞 128프레임. 2초 윈도우 mel 통계는 10초와 다름 → baseline MSE가 부풀려져 임계값 산정과 어긋남."),
]

y = Inches(1.7)
for title, body in details_2:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(0.12), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    add_text_box(slide, Inches(0.65), y, Inches(12.3), Inches(0.4),
                 title, font_size=15, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.65), y + Inches(0.4), Inches(12.3), Inches(0.65),
                 body, font_size=11, color=DARK)
    y += Inches(1.1)

# ========== Slide 9: Improvement Direction ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, GRAY_BG)
add_header_band(slide, "5. 개선 방향",
                "취약점별 대응 전략 — 데이터 / 모델 / 손실 / 운영")

imp_blocks = [
    ("데이터·전처리",
     ["다중 ID·다중 SNR 학습으로 도메인 일반화",
      "절대 dB 정규화로 변경",
      "SpecAugment / 노이즈 인젝션 등 증강",
      "슬라이딩 윈도우로 학습-추론 분포 일치"], BLUE),
    ("모델 아키텍처",
     ["VAE / Memory-AE / Normalizing Flow",
      "BatchNorm → LayerNorm/InstanceNorm",
      "Skip-connection 최소화로 복원 제한",
      "잠재공간 차원 축소 또는 정규화 강화"], NAVY),
    ("손실 함수",
     ["SSIM / Perceptual Loss 도입",
      "주파수 대역 가중 MSE",
      "Patch-wise Max Error",
      "Mahalanobis 거리 기반 점수"], ACCENT),
    ("운영·임계값",
     ["적응형(Adaptive) 임계값 갱신",
      "Concept Drift 모니터링",
      "다중 임계값 (경고 / 위험)",
      "주기적 재학습 (Online Learning)"], RED),
]

for i, (title, items, color) in enumerate(imp_blocks):
    r = i // 2
    c = i % 2
    left = Inches(0.5 + c * 6.3)
    top = Inches(1.9 + r * 2.6)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.0), Inches(2.4))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(2.0)
    box.shadow.inherit = False
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(6.0), Inches(0.5))
    head.fill.solid(); head.fill.fore_color.rgb = color; head.line.fill.background()
    add_text_box(slide, left + Inches(0.2), top + Inches(0.05), Inches(5.8), Inches(0.4),
                 title, font_size=15, bold=True, color=LIGHT)
    add_bullet_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.7), Inches(1.7),
                   items, font_size=12, color=DARK)

# ========== Slide 10: Conclusion ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_text_box(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
             "결  론  (Conclusion)", font_size=30, bold=True, color=LIGHT)
accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2),
                                 Inches(2.0), Inches(0.06))
accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()

conc_items = [
    "본 프로젝트는 MIMII 데이터셋(0dB valve id_02)으로 CNN 오토인코더 기반 음향 이상탐지 파이프라인을 성공적으로 구축함.",
    "정상음만 학습하는 비지도 접근으로 결함 라벨 부족 문제를 우회하였고, 슬라이딩 윈도우 기반 실시간 모니터링까지 시연.",
    "그러나 단일 ID·고정 임계값·MSE 손실·BatchNorm 의존 등 10가지 구조적 취약점이 존재하여, 실제 산업 현장 배포 시 오탐과 미탐이 빈번할 수 있음.",
    "특히 Generalization Gap, 학습-추론 분포 불일치, Concept Drift는 본 모델이 즉시 직면할 핵심 리스크임.",
    "후속 연구로 VAE / Memory-AE 도입, 다중 도메인 학습, 적응형 임계값, SSIM·Patch-wise 손실로 강건성을 확보할 필요가 있음.",
]
add_bullet_box(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(4.5),
               conc_items, font_size=15, color=LIGHT)

add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             "Thank you for your attention.",
             font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


out_path = "결과보고서_MIMII_CNN_AE_취약점분석.pptx"
prs.save(out_path)
print(f"[OK] PPT saved: {out_path}")
print(f"     slides: {len(prs.slides)}")
